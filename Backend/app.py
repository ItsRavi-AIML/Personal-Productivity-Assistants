from flask import Flask, request, jsonify, session, url_for, redirect
from flask_session import Session
from authlib.integrations.flask_client import OAuth
from flask_cors import CORS
from dotenv import load_dotenv
import os
import PyPDF2
from docx import Document
from pptx import Presentation
import chromadb
from chromadb.utils import embedding_functions
from groq import Groq
os.environ['OAUTHLIB_INSECURE_TRANSPORT'] = '1'

load_dotenv()

app = Flask(__name__)
CORS(app, origins=["http://127.0.0.1:5500", "http://localhost:5500", "http://127.0.0.1:8080", "http://localhost:8080"], supports_credentials=True)

# ── Session config ──
app.secret_key = os.getenv("SECRET_KEY", "nexus2024xyzfixed")
app.config['SESSION_TYPE'] = 'filesystem'
app.config['SESSION_COOKIE_SAMESITE'] = 'Lax'
app.config['SESSION_COOKIE_SECURE'] = False
app.config['SESSION_COOKIE_HTTPONLY'] = True
app.config['SESSION_FILE_DIR'] = './flask_sessions'
app.config['SESSION_PERMANENT'] = False
Session(app)
oauth = OAuth(app)

google = oauth.register(
    name='google',
    client_id=os.getenv("GOOGLE_CLIENT_ID"),
    client_secret=os.getenv("GOOGLE_CLIENT_SECRET"),
    server_metadata_url='https://accounts.google.com/.well-known/openid-configuration',
    client_kwargs={'scope': 'openid email profile'}
)

UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

ALLOWED_EXTENSIONS = {'pdf', 'docx', 'pptx', 'txt', 'md'}

# ── Setup ChromaDB ──
chroma_client = chromadb.PersistentClient(path="./chromadb_store")
embedding_fn = embedding_functions.SentenceTransformerEmbeddingFunction(
    model_name="all-MiniLM-L6-v2"
)
collection = chroma_client.get_or_create_collection(
    name="nexus_docs",
    embedding_function=embedding_fn
)

# ── Setup Groq ──
groq_client = Groq(api_key=os.getenv("GROQ_API_KEY"))

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

# ── Extract text ──
def extract_text(filepath):
    ext = filepath.rsplit('.', 1)[1].lower()
    text = ""

    if ext == "pdf":
        with open(filepath, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() or ""

    elif ext == "docx":
        doc = Document(filepath)
        for para in doc.paragraphs:
            text += para.text + "\n"

    elif ext == "pptx":
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    elif ext in ["txt", "md"]:
        with open(filepath, "r", encoding="utf-8") as f:
            text = f.read()

    return text.strip()

# ── Chunk text ──
def chunk_text(text, chunk_size=500, overlap=50):
    words = text.split()
    chunks = []
    i = 0
    while i < len(words):
        chunk = " ".join(words[i:i+chunk_size])
        chunks.append(chunk)
        i += chunk_size - overlap
    return chunks

# ── Upload Route ──
@app.route("/upload", methods=["POST"])
def upload_file():
    if 'file' not in request.files:
        return jsonify({"error": "No file provided"}), 400

    file = request.files['file']

    if file.filename == '':
        return jsonify({"error": "No file selected"}), 400

    if not allowed_file(file.filename):
        return jsonify({"error": "File type not supported"}), 400

    filepath = os.path.join(UPLOAD_FOLDER, file.filename)
    file.save(filepath)

    text = extract_text(filepath)
    if not text:
        return jsonify({"error": "Could not extract text from file"}), 400

    chunks = chunk_text(text)

    for i, chunk in enumerate(chunks):
        chunk_id = f"{file.filename}_chunk_{i}"
        collection.upsert(
            documents=[chunk],
            ids=[chunk_id],
            metadatas=[{"filename": file.filename, "chunk_index": i}]
        )

    return jsonify({
        "message": "File uploaded and stored successfully",
        "filename": file.filename,
        "chunks_stored": len(chunks)
    })

# ── Query Route ──
@app.route("/query", methods=["POST"])
def query():
    try:
        data = request.get_json(force=True, silent=True)
        
        if not data:
            return jsonify({"error": "Invalid JSON"}), 400

        question = data.get("question", "").strip()

        if not question:
            return jsonify({"error": "No question provided"}), 400

        results = collection.query(
            query_texts=[question],
            n_results=5
        )

        chunks = results["documents"][0]
        context = "\n\n".join(chunks)

        if not context:
            return jsonify({"error": "No relevant documents found"}), 404

        response = groq_client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": "You are NEXUS, a personal productivity assistant. Answer the user's question using only the provided document context. If the answer is not in the context, say so clearly."
                },
                {
                    "role": "user",
                    "content": f"Context from documents:\n{context}\n\nQuestion: {question}"
                }
            ]
        )

        answer = response.choices[0].message.content

        return jsonify({
            "answer": answer,
            "sources": results["metadatas"][0]
        })

    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ── Google Login Route ──
@app.route('/login/google')
def google_login():
    # This generates the URL for the 'google_callback' route
    redirect_uri = url_for('google_callback', _external=True)
    return google.authorize_redirect(redirect_uri)
# ── Google Auth Routes ──
@app.route('/auth/callback')
def google_callback():
    token = google.authorize_access_token()
    user_info = google.get('https://www.googleapis.com/oauth2/v3/userinfo').json()
    
    session['user'] = user_info
    print(f"✅ User Logged In: {user_info['name']}")
    
    name = user_info.get('name', '')
    email = user_info.get('email', '')
    picture = user_info.get('picture', '')
    
    return redirect(f"http://127.0.0.1:5500/Frontend/pages/dashboard.html?name={name}&email={email}&picture={picture}")
# ── NEW: Bridge to send user data to Frontend ──
@app.route("/me")
def get_current_user():
    user = session.get('user')
    if not user:
        return jsonify({"error": "Not logged in"}), 401
    return jsonify(user)


# ── Home ──
@app.route("/")
def home():
    return jsonify({"message": "NEXUS Backend is running"})

# ── List uploaded files ──
@app.route("/files", methods=["GET"])
def list_files():
    files = []
    for filename in os.listdir(UPLOAD_FOLDER):
        filepath = os.path.join(UPLOAD_FOLDER, filename)
        files.append({
            "filename": filename,
            "size": os.path.getsize(filepath)
        })
    return jsonify({"files": files})
# ── Delete file and its chunks ──
@app.route("/files/<filename>", methods=["DELETE"])
def delete_file(filename):
    try:
        filepath = os.path.join(UPLOAD_FOLDER, filename)
        if os.path.exists(filepath):
            os.remove(filepath)
            collection.delete(where={"filename": filename})
            return jsonify({"message": "File deleted successfully"})
        else:
            return jsonify({"error": "File not found"}), 404
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# ── Summarize Route ──
@app.route("/summarize", methods=["POST"])
def summarize():
    try:
        data = request.get_json(force=True, silent=True)
        if not data:
            return jsonify({"error": "Invalid JSON"}), 400

        filename = data.get("filename", "").strip()
        mode = data.get("mode", "deep").strip()

        if not filename:
            return jsonify({"error": "No filename provided"}), 400

        # Get all chunks for this file from ChromaDB
        results = collection.get(
            where={"filename": filename}
        )

        if not results["documents"]:
            return jsonify({"error": "File not found in database"}), 404

        context = "\n\n".join(results["documents"])

        # Mode-based instructions
        mode_prompts = {
            "quick": "Give a short 3-5 bullet point summary of the key points only.",
            "deep": "Give a detailed structured summary with sections: Key Points, Architecture/Main Ideas, Action Items, and Conclusions.",
            "custom": "Summarize focusing on deadlines, action items, and important facts."
        }

        prompt = mode_prompts.get(mode, mode_prompts["deep"])

        response = groq_client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": f"You are NEXUS, a document summarizer. {prompt} Be concise and accurate. Only use information from the provided document."
                },
                {
                    "role": "user",
                    "content": f"Document content:\n{context}\n\nPlease summarize this document."
                }
            ]
        )

        summary = response.choices[0].message.content

        return jsonify({
            "summary": summary,
            "filename": filename,
            "mode": mode,
            "chunks_used": len(results["documents"])
        })

    except Exception as e:
        return jsonify({"error": str(e)}), 500
    
    @app.route("/files/<filename>", methods=["DELETE"])
def delete_file(filename):
    try:
        filepath = os.path.join(UPLOAD_FOLDER, filename)
        if os.path.exists(filepath):
            os.remove(filepath)
            collection.delete(where={"filename": filename})
            return jsonify({"message": "File deleted successfully"})
        else:
            return jsonify({"error": "File not found"}), 404
    except Exception as e:
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    app.run(debug=True, port=5000, use_reloader=False)